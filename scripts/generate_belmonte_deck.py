import os
import random
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageFont
try:
    import qrcode
    from qrcode.image.styledpil import StyledPilImage
    from qrcode.image.styles.moduledrawers import RoundedModuleDrawer
    from qrcode.image.styles.colormasks import RadialGradiantColorMask
except Exception:  # Fallback if styled features are unavailable
    import qrcode

PRIMARY = RGBColor(0x2C, 0x4D, 0x7B)      # Azul oceano
ACCENT = RGBColor(0x4C, 0x9E, 0xD9)       # Azul claro
GOLD = RGBColor(0xF3, 0xA6, 0x4D)         # Dourado pôr do sol
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# PIL tuples for image generation
PRIMARY_PIL = (0x2C, 0x4D, 0x7B)
ACCENT_PIL = (0x4C, 0x9E, 0xD9)
GOLD_PIL = (0xF3, 0xA6, 0x4D)
WHITE_PIL = (0xFF, 0xFF, 0xFF)


def add_title_slide(prs: Presentation):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Background rectangle in primary color
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY
    bg.line.fill.background()

    # Glassmorphism card
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.0), Inches(1.0), prs.slide_width - Inches(2.0), prs.slide_height - Inches(2.0),
    )
    card.fill.solid()
    # Light translucent white
    card.fill.fore_color.rgb = WHITE
    card.fill.transparency = 0.15
    card.line.color.rgb = ACCENT
    card.line.width = Pt(1.5)

    tx = card.text_frame
    tx.clear()
    p = tx.paragraphs[0]
    run = p.add_run()
    run.text = "Guia de Turismo de Belmonte"
    p.alignment = PP_ALIGN.LEFT
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = PRIMARY
    run.font.name = "Poppins"

    p2 = tx.add_paragraph()
    p2.text = "App com IA Bel: roteiros inteligentes, mapa offline e experiências locais."
    p2.alignment = PP_ALIGN.LEFT
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(40, 40, 40)
    p2.font.name = "Montserrat"

    # CTA chip
    chip = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(4.7), Inches(3.1), Inches(0.8)
    )
    chip.fill.solid()
    chip.fill.fore_color.rgb = GOLD
    chip.line.fill.background()
    chip.text_frame.text = "Criar meu roteiro"
    chip.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    chip.text_frame.paragraphs[0].runs[0].font.bold = True
    chip.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
    chip.text_frame.paragraphs[0].runs[0].font.name = "Poppins"


def add_agenda_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(8.0), Inches(1.0))
    tf = title.text_frame
    tf.text = "O que você verá"
    tf.paragraphs[0].runs[0].font.size = Pt(32)
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.color.rgb = PRIMARY
    tf.paragraphs[0].runs[0].font.name = "Poppins"

    bullets = [
        "Visão geral do app e identidade visual",
        "IA Bel: a guia que personaliza sua viagem",
        "Telas-chave: Mapa, Roteiros, Gastronomia, Hospedagem, Comércio",
        "Interações: microanimações, parallax, ícones animados",
        "Uso em pitch e apresentação institucional",
    ]
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.0), Inches(4.5))
    t = box.text_frame
    t.clear()
    for idx, item in enumerate(bullets):
        p = t.add_paragraph() if idx > 0 else t.paragraphs[0]
        p.text = item
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.font.name = "Montserrat"


def add_feature_slide(prs: Presentation, title: str, points: list[str], accent: RGBColor):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Accent header band
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    band.fill.solid()
    band.fill.fore_color.rgb = accent
    band.line.fill.background()

    hdr = slide.shapes.add_textbox(Inches(0.6), Inches(0.2), Inches(8.0), Inches(0.8))
    tf = hdr.text_frame
    tf.text = title
    tf.paragraphs[0].runs[0].font.size = Pt(30)
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.color.rgb = WHITE
    tf.paragraphs[0].runs[0].font.name = "Poppins"

    content = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(4.6))
    ct = content.text_frame
    ct.clear()
    for i, pt in enumerate(points):
        p = ct.add_paragraph() if i > 0 else ct.paragraphs[0]
        p.text = pt
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.font.name = "Montserrat"


def add_screens_overview(prs: Presentation):
    add_feature_slide(
        prs,
        "Telas do App",
        [
            "Tela Inicial – boas-vindas da IA Bel, CTA ‘Criar meu roteiro’",
            "Mapa Interativo – pontos turísticos destacados, opção offline",
            "Roteiro Inteligente – sugestões personalizadas por IA",
            "Gastronomia – restaurantes, avaliações, filtros",
            "Hospedagem – hotéis e reservas diretas",
            "Comércio Local – artesanato e produtos típicos",
            "Favoritos – salve lugares e roteiros",
            "Perfil – dados pessoais e histórico",
        ],
        ACCENT,
    )


def add_bel_ai_slide(prs: Presentation):
    add_feature_slide(
        prs,
        "IA Bel – sua guia pessoal",
        [
            "Mascote digital minimalista, tom tropical e amigável",
            "Sugere restaurantes, passeios e cultura local",
            "Balões de fala leves, com glassmorphism",
            "Presença contextual em todas as telas-chave",
        ],
        PRIMARY,
    )


def add_interactions_slide(prs: Presentation):
    add_feature_slide(
        prs,
        "Interações e Motion",
        [
            "Microanimações suaves em botões e transições",
            "Leve parallax em imagens de destaque",
            "CTAs dourados com hover em azul",
            "Ícones flat coloridos com animação sutil",
        ],
        GOLD,
    )


def add_branding_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(8.0), Inches(1.0))
    tf = title.text_frame
    tf.text = "Identidade Visual"
    tf.paragraphs[0].runs[0].font.size = Pt(32)
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.color.rgb = PRIMARY
    tf.paragraphs[0].runs[0].font.name = "Poppins"

    # Palette chips
    chips = [
        ("Azul oceano #2C4D7B", PRIMARY),
        ("Azul claro #4C9ED9", ACCENT),
        ("Dourado pôr do sol #F3A64D", GOLD),
        ("Branco #FFFFFF", WHITE),
    ]
    x = 0.8
    for label, color in chips:
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.7), Inches(2.3), Inches(0.9))
        rect.fill.solid()
        rect.fill.fore_color.rgb = color
        rect.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(x), Inches(2.7), Inches(2.3), Inches(0.5))
        t = tb.text_frame
        t.text = label
        t.paragraphs[0].runs[0].font.size = Pt(12)
        t.paragraphs[0].runs[0].font.color.rgb = RGBColor(40, 40, 40)
        t.paragraphs[0].runs[0].font.name = "Montserrat"
        x += 2.5

    # Typography
    ty = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(8.5), Inches(2.0))
    tty = ty.text_frame
    tty.clear()
    p1 = tty.paragraphs[0]
    p1.text = "Títulos: Poppins Bold"
    p1.font.size = Pt(18)
    p1.font.color.rgb = RGBColor(40, 40, 40)
    p1.font.name = "Montserrat"
    p2 = tty.add_paragraph()
    p2.text = "Textos: Montserrat Regular"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(40, 40, 40)
    p2.font.name = "Montserrat"


def add_closing_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(8.0), Inches(2.5))
    tf = box.text_frame
    tf.text = "Obrigado!"
    tf.paragraphs[0].runs[0].font.size = Pt(40)
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.color.rgb = PRIMARY
    tf.paragraphs[0].runs[0].font.name = "Poppins"
    p = tf.add_paragraph()
    p.text = "Pronto para pitch, apresentação e divulgação."
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(60, 60, 60)
    p.font.name = "Montserrat"


def _draw_vertical_gradient(draw: ImageDraw.ImageDraw, w: int, h: int, top_color: tuple[int, int, int], bottom_color: tuple[int, int, int]):
    for y in range(h):
        ratio = y / max(1, h - 1)
        r = int(top_color[0] * (1 - ratio) + bottom_color[0] * ratio)
        g = int(top_color[1] * (1 - ratio) + bottom_color[1] * ratio)
        b = int(top_color[2] * (1 - ratio) + bottom_color[2] * ratio)
        draw.line([(0, y), (w, y)], fill=(r, g, b))


def _draw_tropical_blobs(draw: ImageDraw.ImageDraw, w: int, h: int):
    # Abstract tropical blobs using brand colors
    for _ in range(6):
        cx = random.randint(int(w * 0.05), int(w * 0.95))
        cy = random.randint(int(h * 0.05), int(h * 0.95))
        rx = random.randint(int(w * 0.06), int(w * 0.18))
        ry = random.randint(int(h * 0.04), int(h * 0.12))
        color = random.choice([
            (76, 158, 217, 80),  # ACCENT alpha
            (243, 166, 77, 90),  # GOLD alpha
            (255, 255, 255, 50),
        ])
        bbox = [cx - rx, cy - ry, cx + rx, cy + ry]
        draw.ellipse(bbox, fill=color)


def _draw_status_bar(draw: ImageDraw.ImageDraw, w: int, font: ImageFont.ImageFont):
    # Simple top status/header bar
    bar_h = 120
    draw.rectangle([0, 0, w, bar_h], fill=(255, 255, 255, 180))
    # App title
    title = "Belmonte"
    tw = draw.textlength(title, font=font)
    draw.text(((w - tw) / 2, 36), title, font=font, fill=(44, 77, 123))
    # IA Bel chip
    draw.rounded_rectangle([w - 280, 28, w - 40, 92], radius=24, fill=(243, 166, 77))
    draw.text((w - 260, 42), "IA Bel", font=font, fill=(255, 255, 255))


def _ensure_font() -> ImageFont.ImageFont:
    try:
        return ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 48)
    except Exception:
        return ImageFont.load_default()


def _save(img: Image.Image, path: str):
    os.makedirs(os.path.dirname(path), exist_ok=True)
    img.save(path, format="PNG")


def _mock_home(w: int, h: int) -> Image.Image:
    img = Image.new("RGBA", (w, h))
    draw = ImageDraw.Draw(img, "RGBA")
    _draw_vertical_gradient(draw, w, h, ACCENT_PIL, PRIMARY_PIL)
    font = _ensure_font()
    _draw_status_bar(draw, w, font)
    # Hero card
    draw.rounded_rectangle([80, 240, w - 80, 680], radius=48, fill=(255, 255, 255, 210))
    draw.text((120, 300), "Bem-vindo!", font=font, fill=(44, 77, 123))
    # CTA button
    draw.rounded_rectangle([80, 760, w - 80, 900], radius=36, fill=(243, 166, 77))
    draw.text((w / 2 - 280, 800), "Criar meu roteiro", font=font, fill=(255, 255, 255))
    return img


def _mock_map(w: int, h: int) -> Image.Image:
    img = Image.new("RGBA", (w, h))
    draw = ImageDraw.Draw(img, "RGBA")
    _draw_vertical_gradient(draw, w, h, (230, 245, 255), (200, 230, 255))
    font = _ensure_font()
    _draw_status_bar(draw, w, font)
    # Map area
    draw.rounded_rectangle([40, 160, w - 40, h - 180], radius=28, fill=(245, 250, 255, 230))
    # Pins
    for _ in range(14):
        x = random.randint(80, w - 80)
        y = random.randint(220, h - 240)
        draw.ellipse([x - 16, y - 16, x + 16, y + 16], fill=(243, 166, 77))
        draw.rectangle([x - 3, y + 16, x + 3, y + 36], fill=(44, 77, 123))
    return img


def _mock_ai_itinerary(w: int, h: int) -> Image.Image:
    img = Image.new("RGBA", (w, h))
    draw = ImageDraw.Draw(img, "RGBA")
    _draw_vertical_gradient(draw, w, h, (250, 250, 255), (235, 245, 255))
    font = _ensure_font()
    _draw_status_bar(draw, w, font)
    # Chat bubbles
    y = 220
    for i in range(6):
        left = 60 if i % 2 == 0 else w - 760
        right = left + 700
        color = (76, 158, 217, 180) if i % 2 == 0 else (255, 255, 255, 220)
        draw.rounded_rectangle([left, y, right, y + 120], radius=28, fill=color)
        y += 160
    # Suggestion chips
    for i, txt in enumerate(["Praias", "Cultura", "Gastronomia", "Família"]):
        draw.rounded_rectangle([60 + i * 260, h - 260, 280 + i * 260, h - 190], radius=28, fill=(243, 166, 77))
    return img


def _mock_list(w: int, h: int, title_color=(44, 77, 123)) -> Image.Image:
    img = Image.new("RGBA", (w, h))
    draw = ImageDraw.Draw(img, "RGBA")
    _draw_vertical_gradient(draw, w, h, (250, 252, 255), (236, 244, 252))
    font = _ensure_font()
    _draw_status_bar(draw, w, font)
    y = 200
    for _ in range(6):
        draw.rounded_rectangle([60, y, w - 60, y + 220], radius=28, fill=(255, 255, 255, 230))
        # Thumbnail
        draw.rounded_rectangle([80, y + 20, 280, y + 200], radius=16, fill=(76, 158, 217))
        # Title bar
        draw.rectangle([300, y + 20, w - 80, y + 80], fill=title_color)
        # Stars
        sx = 300
        for _s in range(5):
            draw.ellipse([sx, y + 100, sx + 40, y + 140], fill=(243, 166, 77))
            sx += 56
        y += 240
    return img


def _mock_grid(w: int, h: int) -> Image.Image:
    img = Image.new("RGBA", (w, h))
    draw = ImageDraw.Draw(img, "RGBA")
    _draw_vertical_gradient(draw, w, h, (250, 250, 250), (235, 240, 245))
    font = _ensure_font()
    _draw_status_bar(draw, w, font)
    # Grid 2 columns
    margin = 60
    pad = 40
    cell_w = (w - margin * 2 - pad) // 2
    cell_h = 320
    x = margin
    y = 220
    for _ in range(8):
        draw.rounded_rectangle([x, y, x + cell_w, y + cell_h], radius=28, fill=(255, 255, 255, 230))
        draw.rounded_rectangle([x + 20, y + 20, x + cell_w - 20, y + 200], radius=16, fill=(76, 158, 217))
        draw.ellipse([x + cell_w - 80, y + 20, x + cell_w - 20, y + 80], fill=(243, 166, 77))
        x += cell_w + pad
        if x + cell_w > w - margin:
            x = margin
            y += cell_h + pad
    return img


def generate_mockups(mock_dir: str) -> dict[str, str]:
    os.makedirs(mock_dir, exist_ok=True)
    w, h = 1080, 1920
    mapping: dict[str, str] = {}
    items = [
        ("Tela Inicial", _mock_home, "tela_inicial.png"),
        ("Mapa Interativo", _mock_map, "mapa.png"),
        ("Roteiro IA", _mock_ai_itinerary, "roteiro_ia.png"),
        ("Gastronomia", _mock_list, "gastronomia.png"),
        ("Hospedagem", _mock_list, "hospedagem.png"),
        ("Comércio Local", _mock_grid, "comercio_local.png"),
        ("Favoritos", _mock_grid, "favoritos.png"),
        ("Perfil", _mock_list, "perfil.png"),
    ]
    for title, fn, filename in items:
        img = fn(w, h)
        path = os.path.join(mock_dir, filename)
        _save(img, path)
        mapping[title] = path
    return mapping


def add_mockup_gallery(prs: Presentation, title: str, mapping: dict[str, str]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(8.0), Inches(1.0))
    tf = header.text_frame
    tf.text = title
    tf.paragraphs[0].runs[0].font.size = Pt(32)
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.color.rgb = PRIMARY
    tf.paragraphs[0].runs[0].font.name = "Poppins"

    # Grid 2x4 thumbnails
    col_w = Inches(4.8)
    row_h = Inches(3.8)
    left0 = Inches(0.8)
    top0 = Inches(1.6)
    i = 0
    for screen_title, img_path in mapping.items():
        row = i // 2
        col = i % 2
        left = left0 + col * col_w
        top = top0 + row * row_h
        slide.shapes.add_picture(img_path, left, top, width=Inches(3.2))
        cap = slide.shapes.add_textbox(left, top + Inches(3.3), Inches(3.2), Inches(0.5))
        ctf = cap.text_frame
        ctf.text = screen_title
        ctf.paragraphs[0].runs[0].font.size = Pt(14)
        ctf.paragraphs[0].runs[0].font.color.rgb = RGBColor(60, 60, 60)
        ctf.paragraphs[0].runs[0].font.name = "Montserrat"
        i += 1


def add_mockup_detail_slides(prs: Presentation, mapping: dict[str, str]):
    for screen_title, img_path in mapping.items():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        header = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(8.0), Inches(1.0))
        tf = header.text_frame
        tf.text = screen_title
        tf.paragraphs[0].runs[0].font.size = Pt(28)
        tf.paragraphs[0].runs[0].font.bold = True
        tf.paragraphs[0].runs[0].font.color.rgb = PRIMARY
        tf.paragraphs[0].runs[0].font.name = "Poppins"
        slide.shapes.add_picture(img_path, Inches(1.2), Inches(1.4), width=Inches(6.5))


def _make_center_logo(size: int = 220) -> Image.Image:
    # Create a simple circular brand mark with GOLD and "Bel"
    img = Image.new("RGBA", (size, size), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img, "RGBA")
    draw.ellipse([0, 0, size, size], fill=(*GOLD_PIL, 255))
    font = _ensure_font()
    text = "Bel"
    tw = draw.textlength(text, font=font)
    draw.text(((size - tw) / 2, size / 2 - 24), text, font=font, fill=(255, 255, 255))
    return img


def generate_qr_code(url: str, out_path: str) -> str:
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    try:
        qr = qrcode.QRCode(
            version=None,
            error_correction=qrcode.constants.ERROR_CORRECT_H,
            box_size=12,
            border=2,
        )
        qr.add_data(url)
        qr.make(fit=True)
        try:
            img = qr.make_image(
                image_factory=StyledPilImage,
                module_drawer=RoundedModuleDrawer(),
                color_mask=RadialGradiantColorMask(center_color=ACCENT_PIL, edge_color=PRIMARY_PIL),
                back_color=(255, 255, 255),
            )
        except Exception:
            img = qr.make_image(fill_color=ACCENT_PIL, back_color="white")

        img = img.convert("RGBA")
        # Paste center logo
        logo = _make_center_logo(220)
        lw, lh = logo.size
        iw, ih = img.size
        img.alpha_composite(logo, (iw // 2 - lw // 2, ih // 2 - lh // 2))
        img.save(out_path)
        return out_path
    except Exception:
        # Minimal fallback
        img = qrcode.make(url)
        img.save(out_path)
        return out_path


def add_qr_slide(prs: Presentation, qr_path: str, url: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(8.0), Inches(1.0))
    tf = header.text_frame
    tf.text = "Acesse o conteúdo"
    tf.paragraphs[0].runs[0].font.size = Pt(32)
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.color.rgb = PRIMARY
    tf.paragraphs[0].runs[0].font.name = "Poppins"

    slide.shapes.add_picture(qr_path, Inches(1.2), Inches(1.6), height=Inches(5.5))

    cap = slide.shapes.add_textbox(Inches(7.2), Inches(2.2), Inches(5.6), Inches(3.6))
    ctf = cap.text_frame
    ctf.clear()
    p1 = ctf.paragraphs[0]
    p1.text = "Escaneie para acessar:"
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = PRIMARY
    p1.font.name = "Poppins"
    p2 = ctf.add_paragraph()
    p2.text = url
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(60, 60, 60)
    p2.font.name = "Montserrat"


def build_presentation(output_dir: str) -> str:
    prs = Presentation()
    prs.slide_height = Inches(7.5)
    prs.slide_width = Inches(13.333)  # 16:9

    add_title_slide(prs)
    add_agenda_slide(prs)
    add_branding_slide(prs)
    add_screens_overview(prs)
    add_bel_ai_slide(prs)
    add_interactions_slide(prs)
    # Generate mockups and insert
    mock_dir = os.path.join(output_dir, "mockups")
    mapping = generate_mockups(mock_dir)
    add_mockup_gallery(prs, "Mockups das Telas", mapping)
    add_mockup_detail_slides(prs, mapping)
    # QR slide
    url = os.environ.get("BELMONTE_QR_URL", "https://example.com/belmonte")
    qr_dir = os.path.join(output_dir, "qr")
    qr_path = os.path.join(qr_dir, "belmonte_qr.png")
    generate_qr_code(url, qr_path)
    add_qr_slide(prs, qr_path, url)
    add_closing_slide(prs)

    os.makedirs(output_dir, exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M")
    pptx_path = os.path.join(output_dir, f"Belmonte_Turismo_Premium_{timestamp}.pptx")
    prs.save(pptx_path)
    return pptx_path


def try_export_pdf(pptx_path: str) -> str | None:
    # Attempt export via LibreOffice headless
    try:
        out_dir = os.path.dirname(pptx_path)
        exit_code = os.system(
            f"libreoffice --headless --convert-to pdf --outdir '{out_dir}' '{pptx_path}' >/dev/null 2>&1"
        )
        if exit_code == 0:
            base = os.path.splitext(os.path.basename(pptx_path))[0]
            pdf_path = os.path.join(out_dir, f"{base}.pdf")
            return pdf_path if os.path.exists(pdf_path) else None
        return None
    except Exception:
        return None


def main():
    output_dir = os.environ.get("OUTPUT_DIR", "/workspace/output")
    pptx_path = build_presentation(output_dir)
    pdf_path = try_export_pdf(pptx_path)
    print(f"PPTX:\t{pptx_path}")
    print(f"PDF:\t{pdf_path if pdf_path else 'PDF export not available'}")


if __name__ == "__main__":
    main()

